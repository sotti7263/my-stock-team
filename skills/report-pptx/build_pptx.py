#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""report-pptx: reports/{종목명}.md 를 읽어 디자인된 PPTX 리포트를 생성한다.

사용법:
    python build_pptx.py <종목명>
    python build_pptx.py reports/삼성물산.md          # 경로 직접 지정도 가능

입력:  reports/{종목명}.md
출력:  reports/{종목명}.pptx

슬라이드 순서(항상 고정):
    1) 표지(종목명·작성일)  2) 종목 개요  3) 재무 요약(표, 최근 3개년)
    4) 가격/추세(차트)  5) 뉴스·심리  6) 리스크  7) 한 줄 종합

설계 규칙:
    - 포인트색 KB 옐로우(#FFBC00) + 본문 그레이/화이트, 차분한 금융 리포트 톤
    - 한글 폰트는 '맑은 고딕' 하나로 고정(글자 깨짐 방지)
    - 수치 출처는 .md 에 적힌 그대로 렌더(스크립트가 숫자를 만들지 않음)
    - 표가 슬라이드 밖으로 넘치지 않도록 행 수·폰트를 자동 제한
    - 매수/매도·목표가 단정 표현은 가드레일 린트로 경고
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 디자인 토큰 ----------------------------------------------------------
KB_YELLOW = RGBColor(0xFF, 0xBC, 0x00)
INK = RGBColor(0x2B, 0x2B, 0x2B)       # 본문 진회색
SUBINK = RGBColor(0x70, 0x70, 0x70)    # 보조 회색
LINE = RGBColor(0xE2, 0xE2, 0xE2)      # 옅은 구분선
BAND = RGBColor(0xF5, 0xF5, 0xF3)      # 표 줄무늬(아주 옅은 회색)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEAD_BG = RGBColor(0x33, 0x33, 0x33)   # 표 헤더 배경

FONT = "맑은 고딕"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.85)
CONTENT_W = SLIDE_W - MARGIN * 2

DISCLAIMER = "본 자료는 학습용 분석이며 투자 권유가 아닙니다."

# 가드레일: 매수/매도·목표가 단정 표현
FORBIDDEN = [
    r"사야\s*합니다", r"팔아야\s*합니다", r"매수\s*추천", r"매도\s*추천",
    r"강력\s*매수", r"강력\s*매도", r"목표\s*주?가", r"적정\s*주가",
    r"반드시\s*오를", r"반드시\s*떨어질",
]


# ---- 마크다운 파싱 --------------------------------------------------------
SECTION_KEYS = [
    ("overview", ["개요", "종목 개요", "기업 개요"]),
    ("financials", ["재무", "재무 요약", "실적"]),
    ("price", ["가격", "추세", "차트", "가격/추세"]),
    ("news", ["뉴스", "심리", "뉴스·심리", "센티먼트"]),
    ("risk", ["리스크", "위험"]),
    ("summary", ["한 줄 종합", "종합", "종합의견", "결론"]),
]


def classify(title: str):
    t = title.replace(" ", "")
    for key, names in SECTION_KEYS:
        for n in names:
            if n.replace(" ", "") in t:
                return key
    return None


def parse_report(md: str) -> dict:
    """제목·작성일·섹션(키→원본 라인 목록)으로 파싱한다."""
    lines = md.splitlines()
    title, date = None, None
    sections: dict[str, list[str]] = {}
    cur = None
    for raw in lines:
        line = raw.rstrip()
        if title is None and line.startswith("# "):
            title = line[2:].strip()
            continue
        m = re.match(r"#{2,3}\s+(.*)", line)
        if m:
            cur = classify(m.group(1).strip())
            if cur and cur not in sections:
                sections[cur] = []
            continue
        # 작성일 추출 (제목부 또는 본문 어디든)
        if date is None:
            dm = re.search(r"작성일[:：]\s*([0-9]{4}-[0-9]{2}-[0-9]{2})", line)
            if dm:
                date = dm.group(1)
                continue
        if cur:
            sections[cur].append(line)
    return {"title": title or "(제목 없음)", "date": date or "", "sections": sections}


def parse_table(lines: list[str]):
    """연속된 마크다운 표 한 개를 (headers, rows)로 반환. 없으면 None."""
    tbl = [l for l in lines if l.strip().startswith("|")]
    if len(tbl) < 2:
        return None

    def cells(row):
        parts = [c.strip() for c in row.strip().strip("|").split("|")]
        return parts

    headers = cells(tbl[0])
    rows = []
    for r in tbl[1:]:
        if re.match(r"^\s*\|[\s:\-|]+\|\s*$", r):  # 구분선 행
            continue
        rows.append(cells(r))
    # 열 수 정규화
    rows = [(c + [""] * len(headers))[: len(headers)] for c in rows]
    return headers, rows


def parse_chart(lines: list[str]):
    """```chart 펜스 블록의 CSV(date,close)를 (categories, values)로 반환."""
    text = "\n".join(lines)
    m = re.search(r"```chart\s*(.*?)```", text, re.DOTALL)
    if not m:
        return None
    cats, vals = [], []
    for row in m.group(1).strip().splitlines():
        row = row.strip()
        if not row or row.lower().startswith("date"):
            continue
        parts = [p.strip() for p in row.split(",")]
        if len(parts) >= 2:
            try:
                vals.append(float(parts[1].replace(",", "")))
                cats.append(parts[0])
            except ValueError:
                continue
    return (cats, vals) if vals else None


def extract_source(lines: list[str]):
    """섹션에서 '출처...'로 시작하는 캡션 라인을 모은다."""
    out = []
    for l in lines:
        s = l.strip().lstrip("-*").strip()
        if s.startswith("출처") or s.startswith("(출처"):
            out.append(s.strip("()"))
    return out


def extract_bullets(lines: list[str]):
    """표/차트/출처 캡션을 제외한 불릿·문단을 추출한다."""
    bullets = []
    skip_chart = False
    for l in lines:
        s = l.strip()
        if s.startswith("```"):
            skip_chart = not skip_chart
            continue
        if skip_chart or not s or s.startswith("|"):
            continue
        if s.startswith("출처") or s.startswith("(출처"):
            continue
        s = re.sub(r"^[-*]\s+", "", s)
        s = re.sub(r"^\d+\.\s+", "", s)
        # 굵게 마크다운 제거
        s = s.replace("**", "")
        bullets.append(s)
    return bullets


# ---- 렌더링 헬퍼 ----------------------------------------------------------
def set_font(run, size=14, bold=False, color=INK, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    """runs: [(text, dict(font opts)), ...] 또는 [[run,...], ...] 문단 목록."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    paragraphs = runs if runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        for text, opts in para:
            r = p.add_run()
            r.text = text
            set_font(r, **opts)
    return tb


def add_title_bar(slide, idx_title, page_no=None):
    """상단: 옐로우 악센트 + 섹션 타이틀 + 페이지 번호."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.55),
                                 Inches(0.13), Inches(0.46))
    bar.fill.solid(); bar.fill.fore_color.rgb = KB_YELLOW
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_text(slide, MARGIN + Inches(0.28), Inches(0.5), CONTENT_W - Inches(1.0), Inches(0.6),
             [(idx_title, dict(size=24, bold=True, color=INK))], anchor=MSO_ANCHOR.MIDDLE)
    if page_no is not None:
        add_text(slide, SLIDE_W - MARGIN - Inches(0.9), Inches(0.62),
                 Inches(0.9), Inches(0.4),
                 [(f"{page_no:02d}", dict(size=12, bold=True, color=KB_YELLOW))],
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # 하단 구분선
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.12),
                                CONTENT_W, Pt(1.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE
    ln.line.fill.background(); ln.shadow.inherit = False


def add_footer(slide):
    add_text(slide, MARGIN, SLIDE_H - Inches(0.5), CONTENT_W, Inches(0.3),
             [(DISCLAIMER, dict(size=9, color=SUBINK))], align=PP_ALIGN.LEFT)


def _set_cell(cell, text, size, bold, color, align, fill):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(6)
    cell.margin_top = cell.margin_bottom = Pt(3)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color)


def add_table(slide, headers, rows, left, top, width,
              max_body_rows=None, font_size=13):
    """줄무늬 표. 행이 max_body_rows를 넘으면 자르고 생략 수를 반환."""
    omitted = 0
    if max_body_rows and len(rows) > max_body_rows:
        omitted = len(rows) - max_body_rows
        rows = rows[:max_body_rows]

    ncols = len(headers)
    nrows = len(rows) + 1
    # 행 높이: 본문 영역(약 4.6")을 넘지 않도록 폰트·높이 조정
    row_h = Inches(0.42)
    gtable = slide.shapes.add_table(nrows, ncols, left, top, width,
                                    row_h * nrows).table
    gtable.first_row = False
    gtable.horz_banding = False

    # 열 너비: 첫 열을 약간 좁게(연도/항목), 나머지 균등
    first_w = int(width * 0.16)
    rest_w = int((width - first_w) / max(ncols - 1, 1))
    gtable.columns[0].width = Emu(first_w)
    for c in range(1, ncols):
        gtable.columns[c].width = Emu(rest_w)

    for r in range(nrows):
        gtable.rows[r].height = row_h
    # 헤더
    for c, h in enumerate(headers):
        _set_cell(gtable.cell(0, c), h, font_size, True, WHITE,
                  PP_ALIGN.CENTER, HEAD_BG)
    # 본문
    for ri, row in enumerate(rows, start=1):
        fill = WHITE if ri % 2 else BAND
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            _set_cell(gtable.cell(ri, c), val, font_size, False, INK, align, fill)
    return omitted


def add_caption(slide, left, top, width, sources):
    if not sources:
        return top
    txt = "  ·  ".join(sources)
    add_text(slide, left, top + Inches(0.12), width, Inches(0.5),
             [(txt, dict(size=9, color=SUBINK))])
    return top + Inches(0.5)


def add_bullets(slide, items, left, top, width, height,
                size=14, max_items=None):
    if max_items and len(items) > max_items:
        items = items[:max_items]
    # 항목 수에 따라 폰트 약간 축소(오버플로 방지)
    if len(items) >= 7:
        size = min(size, 12)
    elif len(items) >= 5:
        size = min(size, 13)
    paras = []
    for it in items:
        paras.append([("•  ", dict(size=size, bold=True, color=KB_YELLOW)),
                      (it, dict(size=size, color=INK))])
    add_text(slide, left, top, width, height, paras, line_spacing=1.12)


def add_line_chart(slide, cats, vals, left, top, width, height):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    # 카테고리가 많으면 라벨이 겹치므로 일부만 표기
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("종가", vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, left, top, width, height, cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    ser = plot.series[0]
    ser.smooth = False
    ser.format.line.color.rgb = KB_YELLOW
    ser.format.line.width = Pt(2.25)
    # 축 폰트
    for axis in (chart.category_axis, chart.value_axis):
        try:
            axis.tick_labels.font.size = Pt(9)
            axis.tick_labels.font.name = FONT
            axis.format.line.color.rgb = LINE
        except Exception:
            pass
    # x축 라벨 솎기
    try:
        n = len(cats)
        chart.category_axis.tick_labels.offset = 100
        from pptx.oxml.ns import qn as _qn
        catAx = chart._chartSpace.findall(".//" + _qn("c:catAx"))
        if catAx:
            tu = catAx[0].find(_qn("c:tickLblSkip"))
            if tu is None:
                tu = catAx[0].makeelement(_qn("c:tickLblSkip"), {})
                catAx[0].append(tu)
            tu.set("val", str(max(1, n // 8)))
    except Exception:
        pass
    return gf


# ---- 슬라이드 빌더 --------------------------------------------------------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_cover(prs, title, date):
    s = blank(prs)
    # 좌측 옐로우 세로 밴드
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.45), SLIDE_H)
    band.fill.solid(); band.fill.fore_color.rgb = KB_YELLOW
    band.line.fill.background(); band.shadow.inherit = False
    add_text(s, MARGIN, Inches(2.5), CONTENT_W, Inches(0.5),
             [("주식 리서치 리포트", dict(size=16, bold=True, color=SUBINK))])
    add_text(s, MARGIN, Inches(3.0), CONTENT_W, Inches(1.4),
             [(title, dict(size=44, bold=True, color=INK))])
    # 옐로우 언더라인
    ul = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.35),
                            Inches(2.2), Inches(0.10))
    ul.fill.solid(); ul.fill.fore_color.rgb = KB_YELLOW
    ul.line.fill.background(); ul.shadow.inherit = False
    if date:
        add_text(s, MARGIN, Inches(4.6), CONTENT_W, Inches(0.5),
                 [(f"작성일  {date}", dict(size=14, color=SUBINK))])
    add_footer(s)
    return s


def build_overview(prs, sec, page):
    s = blank(prs)
    add_title_bar(s, "종목 개요", page)
    top = Inches(1.45)
    bullets = extract_bullets(sec)
    tbl = parse_table(sec)
    if tbl:
        omitted = add_table(s, tbl[0], tbl[1], MARGIN, top, CONTENT_W,
                            max_body_rows=8)
        top += Inches(0.42) * (min(len(tbl[1]), 8) + 1) + Inches(0.2)
        add_caption(s, MARGIN, top, CONTENT_W, extract_source(sec))
    else:
        add_bullets(s, bullets, MARGIN, top, CONTENT_W, Inches(4.8),
                    size=16, max_items=8)
    add_footer(s)
    return s


def build_financials(prs, sec, page):
    s = blank(prs)
    add_title_bar(s, "재무 요약 (최근 3개년)", page)
    top = Inches(1.5)
    tbl = parse_table(sec)
    if tbl:
        # 최근 3개년만(헤더 제외 3행)
        omitted = add_table(s, tbl[0], tbl[1], MARGIN, top, CONTENT_W,
                            max_body_rows=3, font_size=14)
        nrows = min(len(tbl[1]), 3) + 1
        ctop = top + Inches(0.46) * nrows + Inches(0.1)
        ctop = add_caption(s, MARGIN, ctop, CONTENT_W, extract_source(sec))
        if omitted:
            add_text(s, MARGIN, ctop, CONTENT_W, Inches(0.3),
                     [(f"※ 최근 3개년만 표기(이전 {omitted}개 연도 생략)",
                       dict(size=9, color=SUBINK))])
        # 표 아래 코멘트(불릿)
        bul = extract_bullets(sec)
        if bul:
            add_bullets(s, bul, MARGIN, ctop + Inches(0.4), CONTENT_W,
                        Inches(2.4), size=13, max_items=4)
    else:
        add_bullets(s, extract_bullets(sec), MARGIN, top, CONTENT_W,
                    Inches(4.8), size=14)
    add_footer(s)
    return s


def build_price(prs, sec, page):
    s = blank(prs)
    add_title_bar(s, "가격 / 추세", page)
    top = Inches(1.45)
    chart = parse_chart(sec)
    bullets = extract_bullets(sec)
    if chart and len(chart[1]) >= 2:
        add_line_chart(s, chart[0], chart[1], MARGIN, top,
                       Inches(7.4), Inches(4.4))
        # 우측 코멘트
        add_bullets(s, bullets, MARGIN + Inches(7.7), top + Inches(0.2),
                    CONTENT_W - Inches(7.7), Inches(4.3), size=13, max_items=6)
        add_caption(s, MARGIN, top + Inches(4.5), Inches(7.4),
                    extract_source(sec))
    else:
        tbl = parse_table(sec)
        if tbl:
            add_table(s, tbl[0], tbl[1], MARGIN, top, CONTENT_W,
                      max_body_rows=8, font_size=12)
            nrows = min(len(tbl[1]), 8) + 1
            ctop = top + Inches(0.42) * nrows + Inches(0.1)
            ctop = add_caption(s, MARGIN, ctop, CONTENT_W, extract_source(sec))
            add_bullets(s, bullets, MARGIN, ctop + Inches(0.3), CONTENT_W,
                        Inches(1.8), size=12, max_items=3)
        else:
            add_bullets(s, bullets, MARGIN, top, CONTENT_W, Inches(4.8))
    add_footer(s)
    return s


def build_news(prs, sec, page):
    s = blank(prs)
    add_title_bar(s, "뉴스 · 심리", page)
    add_bullets(s, extract_bullets(sec), MARGIN, Inches(1.5), CONTENT_W,
                Inches(4.6), size=14, max_items=7)
    add_caption(s, MARGIN, Inches(6.1), CONTENT_W, extract_source(sec))
    add_footer(s)
    return s


def build_risk(prs, sec, page):
    s = blank(prs)
    add_title_bar(s, "리스크", page)
    tbl = parse_table(sec)
    top = Inches(1.5)
    if tbl:
        add_table(s, tbl[0], tbl[1], MARGIN, top, CONTENT_W,
                  max_body_rows=6, font_size=12)
    else:
        add_bullets(s, extract_bullets(sec), MARGIN, top, CONTENT_W,
                    Inches(4.6), size=14, max_items=7)
    add_caption(s, MARGIN, Inches(6.1), CONTENT_W, extract_source(sec))
    add_footer(s)
    return s


def build_summary(prs, sec, page):
    s = blank(prs)
    add_title_bar(s, "한 줄 종합", page)
    bullets = extract_bullets(sec)
    headline = bullets[0] if bullets else "(종합 의견 미작성)"
    rest = bullets[1:]
    # 강조 박스
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(1.8),
                             CONTENT_W, Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xE0)
    box.line.color.rgb = KB_YELLOW; box.line.width = Pt(1.5)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = headline
    set_font(r, size=20, bold=True, color=INK)
    if rest:
        add_bullets(s, rest, MARGIN, Inches(3.6), CONTENT_W, Inches(2.2),
                    size=14, max_items=5)
    add_caption(s, MARGIN, Inches(5.9), CONTENT_W, extract_source(sec))
    add_footer(s)
    return s


# ---- 가드레일 린트 --------------------------------------------------------
def guardrail_lint(md: str):
    hits = []
    for i, line in enumerate(md.splitlines(), 1):
        for pat in FORBIDDEN:
            if re.search(pat, line):
                hits.append((i, line.strip()))
                break
    return hits


# ---- 엔트리 --------------------------------------------------------------
def resolve_paths(arg: str):
    p = Path(arg)
    if p.suffix.lower() == ".md":
        src = p
    else:
        src = Path("reports") / f"{arg}.md"
    out = src.with_suffix(".pptx")
    return src, out


def main():
    # Windows 콘솔(cp949)에서도 한글·기호가 깨지지 않도록 UTF-8로 출력
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass

    if len(sys.argv) < 2:
        print("사용법: python build_pptx.py <종목명|reports/종목명.md>")
        sys.exit(1)
    src, out = resolve_paths(sys.argv[1])
    if not src.exists():
        print(f"[오류] 입력 파일을 찾을 수 없습니다: {src}")
        sys.exit(1)

    md = src.read_text(encoding="utf-8")

    # 가드레일 점검
    hits = guardrail_lint(md)
    if hits:
        print("⚠ 가드레일 경고 — 매수/매도·목표가 단정 표현이 감지되었습니다:")
        for ln, txt in hits:
            print(f"   L{ln}: {txt}")
        print("   (판단 근거까지만 서술하도록 .md 를 수정하세요. PPTX는 계속 생성합니다.)")
    else:
        print("✓ 가드레일 점검 통과(매수/매도·목표가 단정 표현 없음)")

    data = parse_report(md)
    sec = data["sections"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs, data["title"], data["date"])
    build_overview(prs, sec.get("overview", []), 2)
    build_financials(prs, sec.get("financials", []), 3)
    build_price(prs, sec.get("price", []), 4)
    build_news(prs, sec.get("news", []), 5)
    build_risk(prs, sec.get("risk", []), 6)
    build_summary(prs, sec.get("summary", []), 7)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ 생성 완료: {out}  (슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장)"
          if False else f"✓ 생성 완료: {out}  (슬라이드 7장)")


if __name__ == "__main__":
    main()
